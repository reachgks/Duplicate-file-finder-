import os
import re
import string
import pypdf
import docx
import pptx
import openpyxl

def extract_printable_strings(file_path: str, max_chars: int = 2000) -> str:
    """
    Fallback method for legacy files (.doc, .ppt, .xls).
    Scans the binary file and extracts sequences of printable ASCII/UTF-8 characters.
    Very useful for getting document properties, titles, and headers from binary dumps.
    """
    try:
        with open(file_path, "rb") as f:
            data = f.read(1024 * 128)  # Read first 128KB
        
        # Regex to find sequences of printable characters of length 4 or more
        words = re.findall(rb'[a-zA-Z0-9\s\-_.,()\[\]]{4,}', data)
        text_parts = []
        char_count = 0
        
        for w in words:
            try:
                decoded = w.decode('utf-8', errors='ignore').strip()
                if decoded:
                    text_parts.append(decoded)
                    char_count += len(decoded)
                    if char_count >= max_chars:
                        break
            except Exception:
                continue
        
        return " ".join(text_parts)[:max_chars]
    except Exception as e:
        return f"[Fallback Parsing Error: {str(e)}]"

def parse_pdf(file_path: str, max_chars: int = 4000) -> str:
    """Extract text from a PDF file."""
    try:
        reader = pypdf.PdfReader(file_path)
        text_parts = []
        char_count = 0
        
        # Read first few pages to understand context
        for i, page in enumerate(reader.pages):
            page_text = page.extract_text()
            if page_text:
                text_parts.append(page_text)
                char_count += len(page_text)
                if char_count >= max_chars:
                    break
        
        return "\n".join(text_parts)[:max_chars]
    except Exception as e:
        return f"[PDF Parsing Error: {str(e)}]"

def parse_docx(file_path: str, max_chars: int = 4000) -> str:
    """Extract text from a Word (.docx) file."""
    try:
        doc = docx.Document(file_path)
        text_parts = []
        char_count = 0
        
        # Extract from paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
                char_count += len(para.text)
                if char_count >= max_chars:
                    break
                    
        # Extract from tables if text is still short
        if char_count < max_chars:
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        text_parts.append(row_text)
                        char_count += len(row_text)
                        if char_count >= max_chars:
                            break
                if char_count >= max_chars:
                    break
                    
        return "\n".join(text_parts)[:max_chars]
    except Exception as e:
        return f"[Word Document Parsing Error: {str(e)}]"

def parse_pptx(file_path: str, max_chars: int = 4000) -> str:
    """Extract text from a PowerPoint (.pptx) presentation."""
    try:
        prs = pptx.Presentation(file_path)
        text_parts = []
        char_count = 0
        
        for slide_idx, slide in enumerate(prs.slides):
            slide_text = []
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.text.strip():
                            slide_text.append(paragraph.text.strip())
            
            if slide_text:
                combined_slide = f"--- Slide {slide_idx + 1} ---\n" + "\n".join(slide_text)
                text_parts.append(combined_slide)
                char_count += len(combined_slide)
                if char_count >= max_chars:
                    break
                    
        return "\n\n".join(text_parts)[:max_chars]
    except Exception as e:
        return f"[PowerPoint Parsing Error: {str(e)}]"

def parse_xlsx(file_path: str, max_chars: int = 2000) -> str:
    """Extract summary details and text cells from an Excel (.xlsx) spreadsheet."""
    try:
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        text_parts = []
        
        # Extract sheet names
        text_parts.append(f"Excel File with Sheets: {', '.join(wb.sheetnames)}")
        
        # Read the first sheet's header rows and key rows
        if wb.sheetnames:
            sheet = wb[wb.sheetnames[0]]
            # Read first 15 rows, 10 columns
            rows_data = []
            for r_idx, row in enumerate(sheet.iter_rows(max_row=15, max_col=10, values_only=True)):
                row_vals = [str(cell).strip() for cell in row if cell is not None]
                if row_vals:
                    rows_data.append(f"Row {r_idx+1}: " + " | ".join(row_vals))
            
            if rows_data:
                text_parts.append("Sample Data:\n" + "\n".join(rows_data))
                
        return "\n\n".join(text_parts)[:max_chars]
    except Exception as e:
        return f"[Excel Spreadsheet Parsing Error: {str(e)}]"

def extract_file_content(file_path: str) -> str:
    """
    Main entry point for extracting text content from documents.
    Determines type by file extension and delegates.
    """
    if not os.path.exists(file_path):
        return "[Error: File does not exist]"
        
    _, ext = os.path.splitext(file_path.lower())
    
    # Modern Formats
    if ext == ".pdf":
        return parse_pdf(file_path)
    elif ext == ".docx":
        return parse_docx(file_path)
    elif ext == ".pptx":
        return parse_pptx(file_path)
    elif ext == ".xlsx":
        return parse_xlsx(file_path)
        
    # Legacy Formats
    elif ext in [".doc", ".ppt", ".xls"]:
        return extract_printable_strings(file_path)
        
    else:
        return f"[Unsupported File Format: {ext}]"

if __name__ == "__main__":
    # Quick debug test if run directly
    import sys
    if len(sys.argv) > 1:
        test_file = sys.argv[1]
        print(f"Extracting: {test_file}")
        print("-" * 40)
        print(extract_file_content(test_file))
    else:
        print("Usage: python parser.py <file_path>")
